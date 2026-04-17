import argparse
import json
from datetime import datetime
from pathlib import Path

import matplotlib.pyplot as plt
import openpyxl
import pandas as pd
from docx import Document
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


BASE_DIR = Path(__file__).resolve().parent.parent
INPUT_DIR = BASE_DIR / "data" / "input" / "source_pages"
OUTPUT_BASE_DIR = BASE_DIR / "data" / "output" / "monthly-report"

THEME_BLUE = RGBColor(31, 78, 121)
THEME_DARK = RGBColor(64, 64, 64)
WHITE = RGBColor(255, 255, 255)

SLIDE_W = 8.27
SLIDE_H = 11.69

# =========================
# 프로필 정보: 여기만 바꾸면 4,5페이지에 반영
# =========================
ANALYST_PROFILE = {
    "company": "동양생명",
    "name": "애널리스트 이름",
    "photo": str(BASE_DIR / "data" / "input" / "profile" / "analyst.jpg"),
}

MANAGER_PROFILE = {
    "company": "동양생명",
    "name": "펀드매니저 이름",
    "photo": str(BASE_DIR / "data" / "input" / "profile" / "manager.jpg"),
}

SOURCE_FILES = {
    "cover": INPUT_DIR / "cover" / "(P.1) (변경 O) 표지.xlsx",
    "market_main": INPUT_DIR / "market" / "(P.2-3) (변경 O) 01.주요자산 시장점검.xlsx",
    "market_returns": INPUT_DIR / "market" / "(P.2-3) (변경 O) 별첨_수익률.xlsx",
    "market_charts": INPUT_DIR / "market" / "(P.2-3) (변경 O) 별첨_차트.xlsx",
    "analyst_comment": INPUT_DIR / "comments" / "(P.4) (변경 O) 02.애널리스트 코멘트.docx",
    "manager_comment": INPUT_DIR / "comments" / "(P.5) (변경 O) 03.매니저 코멘트.docx",
    "allocation": INPUT_DIR / "allocation" / "(P.6) (변경 O) 04.투자자별 자산배분 전략.xlsx",
    "funds_1": INPUT_DIR / "funds" / "(P.7) (변경 O) 05.변액보험 추천 펀드(1).xlsx",
    "funds_2": INPUT_DIR / "funds" / "(P.8-9) (변경 O) 05.변액보험 추천 펀드(2).xlsx",
    "lineup_savings": INPUT_DIR / "lineups" / "(P.16) (변경 X) 07.(저축성)대상펀드 Line-up.xlsx",
    "lineup_protection": INPUT_DIR / "lineups" / "(P.17) (변경 X) 08.(보장성)대상펀드 Line-up.xlsx",
    "performance": INPUT_DIR / "performance" / "(P.18) (변경 O) 09.변액펀드 성과현황.xlsx",
    "overseas": INPUT_DIR / "overseas" / "(P.19-21) (변경 O) 10.해외펀드 투자현황.xlsx",
    "managers": INPUT_DIR / "managers" / "(P.23) (변경 O) 12.펀드별 운용사 현황(1).xlsx",
}


def parse_args():
    parser = argparse.ArgumentParser(description="Monthly investment report draft generator")
    parser.add_argument("--month", required=False, default=datetime.now().strftime("%Y-%m"))
    parser.add_argument("--mode", required=False, default="draft", choices=["draft", "final"])
    return parser.parse_args()


def ensure_directories(report_month: str, mode: str) -> dict[str, Path]:
    report_root = OUTPUT_BASE_DIR / report_month / mode
    charts_dir = report_root / "charts"
    logs_dir = report_root / "logs"

    report_root.mkdir(parents=True, exist_ok=True)
    charts_dir.mkdir(parents=True, exist_ok=True)
    logs_dir.mkdir(parents=True, exist_ok=True)

    return {
        "report_root": report_root,
        "charts_dir": charts_dir,
        "logs_dir": logs_dir,
    }


def format_issue_label(report_month: str) -> str:
    try:
        year, month = report_month.split("-")
        return f"{int(year)}년 {int(month)}월"
    except Exception:
        return report_month


def safe_load_workbook(path: Path):
    if not path.exists():
        return None
    try:
        return openpyxl.load_workbook(path, data_only=True)
    except Exception:
        return None


def safe_read_docx(path: Path) -> str:
    if not path.exists():
        return ""
    try:
        doc = Document(path)
        texts = []
        for p in doc.paragraphs:
            txt = (p.text or "").strip()
            if txt:
                texts.append(txt)
        return "\n".join(texts)
    except Exception:
        return ""


def summarize_source_status():
    return [
        {"key": key, "path": str(path), "exists": path.exists()}
        for key, path in SOURCE_FILES.items()
    ]


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text,
    font_size=16,
    bold=False,
    align=PP_ALIGN.LEFT,
    font_color=THEME_DARK,
):
    textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = textbox.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.name = "Malgun Gothic"
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = font_color
    return textbox


def add_footer(slide, report_month: str, page_no: int):
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.45),
        Inches(10.95),
        Inches(7.35),
        Inches(0.02),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = THEME_BLUE
    line.line.color.rgb = THEME_BLUE

    add_textbox(
        slide,
        0.5,
        11.0,
        3.8,
        0.2,
        f"동양생명 월간 투자전략 | {report_month}",
        font_size=8,
        font_color=THEME_DARK,
    )
    add_textbox(
        slide,
        7.1,
        11.0,
        0.5,
        0.2,
        str(page_no),
        font_size=8,
        align=PP_ALIGN.RIGHT,
        font_color=THEME_DARK,
    )


def add_cover_background(slide):
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(SLIDE_W),
        Inches(SLIDE_H),
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(31, 78, 121)
    bg.line.color.rgb = RGBColor(31, 78, 121)


def add_report_header(slide, report_month: str, section_no: str, section_title: str):
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(SLIDE_W),
        Inches(0.95),
    )
    header.fill.solid()
    header.fill.fore_color.rgb = THEME_BLUE
    header.line.color.rgb = THEME_BLUE

    add_textbox(
        slide,
        0.28,
        0.14,
        4.8,
        0.22,
        "동양생명 변액보험 월간 투자전략 GUIDE",
        font_size=12,
        bold=True,
        font_color=WHITE,
    )

    add_textbox(
        slide,
        0.28,
        0.42,
        2.2,
        0.18,
        "모집인교육용 / 대외비",
        font_size=7.5,
        font_color=WHITE,
    )

    pill = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(7.05),
        Inches(0.14),
        Inches(0.95),
        Inches(0.34),
    )
    pill.fill.solid()
    pill.fill.fore_color.rgb = RGBColor(79, 128, 189)
    pill.line.color.rgb = WHITE

    add_textbox(
        slide,
        7.12,
        0.205,
        0.8,
        0.12,
        f"{format_issue_label(report_month)}호",
        font_size=7.3,
        align=PP_ALIGN.CENTER,
        font_color=WHITE,
    )

    add_textbox(
        slide,
        0.55,
        1.1,
        0.9,
        0.55,
        section_no,
        font_size=28,
        bold=True,
        font_color=RGBColor(210, 220, 235),
    )

    add_textbox(
        slide,
        0.55,
        1.6,
        4.2,
        0.35,
        section_title,
        font_size=16,
        bold=True,
        font_color=THEME_BLUE,
    )

    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.55),
        Inches(2.08),
        Inches(7.1),
        Inches(0.025),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = THEME_BLUE
    line.line.color.rgb = THEME_BLUE


def add_chart_comment_box(slide, left, top, width, height, title, body):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(242, 246, 250)
    shape.line.color.rgb = RGBColor(190, 205, 220)

    add_textbox(
        slide,
        left + 0.15,
        top + 0.08,
        width - 0.3,
        0.25,
        title,
        font_size=11,
        bold=True,
        font_color=THEME_BLUE,
    )
    add_textbox(
        slide,
        left + 0.15,
        top + 0.38,
        width - 0.3,
        height - 0.45,
        body,
        font_size=10,
        font_color=THEME_DARK,
    )


def add_summary_comment_box(slide, left, top, width, height, title, body):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(248, 250, 253)
    shape.line.color.rgb = RGBColor(200, 210, 225)

    add_textbox(
        slide,
        left + 0.12,
        top + 0.08,
        width - 0.24,
        0.24,
        title,
        font_size=10,
        bold=True,
        font_color=THEME_BLUE,
    )
    add_textbox(
        slide,
        left + 0.12,
        top + 0.36,
        width - 0.24,
        height - 0.44,
        body,
        font_size=9,
        font_color=THEME_DARK,
    )


def add_profile_box(slide, left, top, width, height, company, name, photo_path):
    outer = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    outer.fill.solid()
    outer.fill.fore_color.rgb = RGBColor(248, 250, 253)
    outer.line.color.rgb = RGBColor(190, 205, 220)

    photo_box_left = left + 0.12
    photo_box_top = top + 0.14
    photo_box_w = 1.15
    photo_box_h = 1.35

    photo_frame = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(photo_box_left),
        Inches(photo_box_top),
        Inches(photo_box_w),
        Inches(photo_box_h),
    )
    photo_frame.fill.solid()
    photo_frame.fill.fore_color.rgb = WHITE
    photo_frame.line.color.rgb = RGBColor(180, 190, 205)

    photo = Path(photo_path) if photo_path else None
    if photo and photo.exists():
        try:
            slide.shapes.add_picture(
                str(photo),
                Inches(photo_box_left + 0.02),
                Inches(photo_box_top + 0.02),
                width=Inches(photo_box_w - 0.04),
                height=Inches(photo_box_h - 0.04),
            )
        except Exception:
            add_textbox(
                slide,
                photo_box_left + 0.12,
                photo_box_top + 0.55,
                0.9,
                0.2,
                "사진",
                font_size=9,
                align=PP_ALIGN.CENTER,
                font_color=THEME_DARK,
            )
    else:
        add_textbox(
            slide,
            photo_box_left + 0.12,
            photo_box_top + 0.55,
            0.9,
            0.2,
            "사진",
            font_size=9,
            align=PP_ALIGN.CENTER,
            font_color=THEME_DARK,
        )

    add_textbox(
        slide,
        left + 1.42,
        top + 0.22,
        width - 1.58,
        0.22,
        "회사명",
        font_size=8.5,
        bold=True,
        font_color=THEME_BLUE,
    )
    add_textbox(
        slide,
        left + 1.42,
        top + 0.46,
        width - 1.58,
        0.22,
        company,
        font_size=10,
        font_color=THEME_DARK,
    )

    add_textbox(
        slide,
        left + 1.42,
        top + 0.86,
        width - 1.58,
        0.22,
        "이름",
        font_size=8.5,
        bold=True,
        font_color=THEME_BLUE,
    )
    add_textbox(
        slide,
        left + 1.42,
        top + 1.10,
        width - 1.58,
        0.26,
        name,
        font_size=11,
        bold=True,
        font_color=THEME_DARK,
    )


def format_table_text(cell, font_size=9, bold=False, align=PP_ALIGN.CENTER, font_color=THEME_DARK):
    tf = cell.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for paragraph in tf.paragraphs:
        paragraph.alignment = align
        for run in paragraph.runs:
            run.font.name = "Malgun Gothic"
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = font_color


def style_table_header(cell):
    cell.fill.solid()
    cell.fill.fore_color.rgb = THEME_BLUE
    format_table_text(cell, font_size=9.5, bold=True, align=PP_ALIGN.CENTER, font_color=WHITE)


def style_table_body(cell, align):
    cell.fill.solid()
    cell.fill.fore_color.rgb = WHITE
    format_table_text(cell, font_size=8.2, bold=False, align=align, font_color=THEME_DARK)


def set_column_widths(table, widths_in_inches):
    for idx, width in enumerate(widths_in_inches):
        table.columns[idx].width = Inches(width)


def find_keyword_positions(ws, keyword: str):
    hits = []
    for row in ws.iter_rows():
        for cell in row:
            val = cell.value
            if val is None:
                continue
            text = str(val).strip()
            if keyword in text:
                hits.append((cell.row, cell.column, text))
    return hits


def row_joined_text(ws, row_idx: int):
    vals = []
    for c in range(1, ws.max_column + 1):
        v = ws.cell(row_idx, c).value
        if v is not None and str(v).strip():
            vals.append(str(v).strip())
    return " ".join(vals).strip()


def extract_block_below_keyword(ws, keyword: str, stop_keywords=None, max_rows=12):
    stop_keywords = stop_keywords or []
    hits = find_keyword_positions(ws, keyword)
    if not hits:
        return ""

    start_row, start_col, _ = hits[-1]
    lines = []

    same_row_vals = []
    for c in range(start_col + 1, min(ws.max_column, start_col + 12) + 1):
        v = ws.cell(start_row, c).value
        if v is not None and str(v).strip():
            same_row_vals.append(str(v).strip())
    same_row_text = " ".join(same_row_vals).strip()
    if same_row_text:
        lines.append(same_row_text)

    for r in range(start_row + 1, min(ws.max_row, start_row + max_rows) + 1):
        joined = row_joined_text(ws, r)
        if not joined:
            continue
        if any(k in joined for k in stop_keywords):
            break
        lines.append(joined)

    return "\n".join(lines).strip()


def normalize_perf_df(df: pd.DataFrame) -> pd.DataFrame:
    if df is None or df.empty:
        return pd.DataFrame()

    df = df.copy()
    df.columns = [str(c).strip() for c in df.columns]

    col_map = {}
    for col in df.columns:
        lc = col.lower()
        if "구분" in col:
            col_map[col] = "구분"
        elif "지수명" in col or "자산명" in col:
            col_map[col] = "자산명"
        elif "1개월" in col or lc == "1m":
            col_map[col] = "1개월"
        elif "3개월" in col or lc == "3m":
            col_map[col] = "3개월"
        elif "6개월" in col or lc == "6m":
            col_map[col] = "6개월"
        elif col in ["1년", "1Y", "1y"]:
            col_map[col] = "1년"
        elif col in ["3년", "3Y", "3y"]:
            col_map[col] = "3년"

    df = df.rename(columns=col_map)

    required = ["자산명", "1개월", "3개월", "6개월", "1년", "3년"]
    if not all(col in df.columns for col in required):
        return pd.DataFrame()

    keep = ["구분", "자산명", "1개월", "3개월", "6개월", "1년", "3년"]
    keep = [c for c in keep if c in df.columns]
    df = df[keep].copy()

    df = df.dropna(subset=["자산명"])
    df["자산명"] = df["자산명"].astype(str).str.strip()
    df = df[df["자산명"] != ""]
    return df.head(12)


def get_market_perf_table_from_returns_file():
    path = SOURCE_FILES["market_returns"]
    if not path.exists():
        return pd.DataFrame()

    try:
        xl = pd.ExcelFile(path)
        for sheet in xl.sheet_names:
            try:
                raw = pd.read_excel(path, sheet_name=sheet)
                normalized = normalize_perf_df(raw)
                if not normalized.empty:
                    return normalized
            except Exception:
                pass
    except Exception:
        pass

    return pd.DataFrame()


def get_market_perf_table_from_main():
    wb = safe_load_workbook(SOURCE_FILES["market_main"])
    if wb is None or "월간투자전략" not in wb.sheetnames:
        return pd.DataFrame()

    ws = wb["월간투자전략"]
    rows = []
    current_group = None
    for r in range(7, 16):
        group_val = ws.cell(r, 3).value
        asset_name = ws.cell(r, 4).value
        if group_val:
            current_group = str(group_val)
        if asset_name:
            rows.append({
                "구분": current_group,
                "자산명": str(asset_name),
                "1개월": ws.cell(r, 8).value,
                "3개월": ws.cell(r, 9).value,
                "6개월": ws.cell(r, 10).value,
                "1년": ws.cell(r, 11).value,
                "3년": ws.cell(r, 12).value,
            })
    return pd.DataFrame(rows)


def get_market_perf_table():
    df = get_market_perf_table_from_returns_file()
    if not df.empty:
        return df
    return get_market_perf_table_from_main()


def get_cover_info():
    wb = safe_load_workbook(SOURCE_FILES["cover"])
    if wb is None or "표지" not in wb.sheetnames:
        return {
            "title": "동양생명 월간 투자전략",
            "subtitle": "변액보험 월간 투자전략 GUIDE",
            "notices": [],
        }

    ws = wb["표지"]
    notices = []
    for r in [31, 32, 34]:
        v = ws.cell(r, 2).value
        if v:
            notices.append(str(v))

    return {
        "title": "동양생명 월간 투자전략",
        "subtitle": "변액보험 월간 투자전략 GUIDE",
        "notices": notices,
    }


def get_market_comments():
    wb = safe_load_workbook(SOURCE_FILES["market_main"])
    if wb is None or "월간투자전략" not in wb.sheetnames:
        return {
            "base_date": "",
            "global_outlook": "",
            "domestic_stocks": "",
            "domestic_bonds": "",
        }

    ws = wb["월간투자전략"]
    base_date = ws.cell(4, 12).value or ""

    global_text = extract_block_below_keyword(
        ws,
        "글로벌 경기 전망",
        stop_keywords=["국내주식", "국내채권", "애널리스트 코멘트", "매니저 코멘트"],
        max_rows=10,
    )
    stock_text = extract_block_below_keyword(
        ws,
        "국내주식",
        stop_keywords=["국내채권", "글로벌 경기 전망", "애널리스트 코멘트", "매니저 코멘트"],
        max_rows=10,
    )
    bond_text = extract_block_below_keyword(
        ws,
        "국내채권",
        stop_keywords=["국내주식", "글로벌 경기 전망", "애널리스트 코멘트", "매니저 코멘트"],
        max_rows=10,
    )

    if not global_text:
        global_text = str(ws.cell(21, 3).value or "")
    if not stock_text:
        stock_text = str(ws.cell(26, 3).value or "")
    if not bond_text:
        bond_text = str(ws.cell(31, 3).value or "")

    return {
        "base_date": str(base_date),
        "global_outlook": str(global_text),
        "domestic_stocks": str(stock_text),
        "domestic_bonds": str(bond_text),
    }


def get_analyst_comment():
    text = safe_read_docx(SOURCE_FILES["analyst_comment"])
    return text if text else "애널리스트 코멘트 파일을 읽지 못했습니다. (.doc 파일이면 .docx로 변환 필요)"


def get_manager_comment():
    text = safe_read_docx(SOURCE_FILES["manager_comment"])
    return text if text else "매니저 코멘트 파일을 읽지 못했습니다."


def get_allocation_data():
    wb = safe_load_workbook(SOURCE_FILES["allocation"])
    if wb is None or "월간투자전략" not in wb.sheetnames:
        return {"reason": "", "current_alloc": pd.DataFrame()}

    ws = wb["월간투자전략"]
    reason = ws.cell(31, 3).value or ""

    rows = []
    for r in range(19, 23):
        asset_name = ws.cell(r, 19).value
        weight = ws.cell(r, 20).value
        if asset_name and weight is not None:
            rows.append({
                "asset_name": str(asset_name),
                "weight": float(weight),
            })

    return {
        "reason": str(reason),
        "current_alloc": pd.DataFrame(rows),
    }


def get_recommended_fund_table():
    wb = safe_load_workbook(SOURCE_FILES["funds_2"])
    if wb is None or "월간투자전략" not in wb.sheetnames:
        return pd.DataFrame()

    ws = wb["월간투자전략"]
    rows = []
    for r in range(7, 17):
        fund_name = ws.cell(r, 4).value
        if fund_name:
            rows.append({
                "유형": ws.cell(r, 3).value,
                "펀드명": fund_name,
                "순자산": ws.cell(r, 6).value,
                "설정일": ws.cell(r, 7).value,
                "펀드등급": ws.cell(r, 8).value,
                "1Y": ws.cell(r, 9).value,
                "2Y": ws.cell(r, 10).value,
                "3Y": ws.cell(r, 11).value,
                "YTD": ws.cell(r, 12).value,
                "설정이후": ws.cell(r, 13).value,
            })
    return pd.DataFrame(rows)


def get_performance_table():
    wb = safe_load_workbook(SOURCE_FILES["performance"])
    if wb is None or "월간투자전략" not in wb.sheetnames:
        return pd.DataFrame()

    ws = wb["월간투자전략"]
    rows = []
    for r in range(17, 29):
        fund_name = ws.cell(r, 4).value
        if fund_name:
            rows.append({
                "자산군": ws.cell(r, 2).value,
                "스타일": ws.cell(r, 3).value,
                "펀드명": fund_name,
                "순자산": ws.cell(r, 5).value,
                "주식비중": ws.cell(r, 6).value,
                "1M": ws.cell(r, 7).value,
                "3M": ws.cell(r, 8).value,
                "6M": ws.cell(r, 9).value,
                "1Y": ws.cell(r, 10).value,
            })
    return pd.DataFrame(rows)


def get_manager_status_table():
    wb = safe_load_workbook(SOURCE_FILES["managers"])
    if wb is None or "운용사현황" not in wb.sheetnames:
        return pd.DataFrame()

    ws = wb["운용사현황"]
    rows = []
    for r in range(4, 14):
        fund_name = ws.cell(r, 3).value
        if fund_name:
            rows.append({
                "펀드명": fund_name,
                "운용사": ws.cell(r, 4).value,
                "수탁회사": ws.cell(r, 5).value,
                "사무관리회사": ws.cell(r, 6).value,
            })
    return pd.DataFrame(rows)


def create_market_perf_summary_chart(perf_df: pd.DataFrame, charts_dir: Path):
    if perf_df.empty or "자산명" not in perf_df.columns or "1년" not in perf_df.columns:
        return None

    plot_df = perf_df.copy()
    plot_df["1년"] = pd.to_numeric(plot_df["1년"], errors="coerce")
    plot_df = plot_df.dropna(subset=["1년"])
    if plot_df.empty:
        return None

    fig, ax = plt.subplots(figsize=(8, 4.4))
    ax.bar(plot_df["자산명"], plot_df["1년"])
    ax.set_title("주요 자산 1년 수익률")
    ax.set_ylabel("수익률(%)")
    plt.xticks(rotation=45, ha="right")

    output_path = charts_dir / "asset_return_1y.png"
    fig.savefig(output_path, bbox_inches="tight", dpi=200)
    plt.close(fig)
    return str(output_path)


def create_market_dual_chart(perf_df: pd.DataFrame, charts_dir: Path):
    if perf_df.empty or "자산명" not in perf_df.columns:
        return None, None

    base_df = perf_df.copy()
    base_df["1개월"] = pd.to_numeric(base_df["1개월"], errors="coerce")
    base_df["3개월"] = pd.to_numeric(base_df["3개월"], errors="coerce")
    base_df["1년"] = pd.to_numeric(base_df["1년"], errors="coerce")
    base_df = base_df.dropna(subset=["1개월", "3개월", "1년"])
    if base_df.empty:
        return None, None

    df1 = base_df.head(6)
    fig1, ax1 = plt.subplots(figsize=(5.8, 3.6))
    ax1.plot(df1["자산명"], df1["1개월"], marker="o", label="1개월")
    ax1.plot(df1["자산명"], df1["3개월"], marker="o", label="3개월")
    ax1.set_title("주요 자산 단기 수익률 비교")
    ax1.legend()
    plt.xticks(rotation=45, ha="right")
    p1 = charts_dir / "market_chart_1.png"
    fig1.savefig(p1, bbox_inches="tight", dpi=200)
    plt.close(fig1)

    df2 = base_df.head(6)
    fig2, ax2 = plt.subplots(figsize=(5.8, 3.6))
    ax2.bar(df2["자산명"], df2["1년"])
    ax2.set_title("주요 자산 1년 수익률")
    plt.xticks(rotation=45, ha="right")
    p2 = charts_dir / "market_chart_2.png"
    fig2.savefig(p2, bbox_inches="tight", dpi=200)
    plt.close(fig2)

    return str(p1), str(p2)


def create_fund_chart(fund_df: pd.DataFrame, charts_dir: Path):
    if fund_df.empty:
        return None

    plot_df = fund_df.copy()
    plot_df["1Y"] = pd.to_numeric(plot_df["1Y"], errors="coerce")
    plot_df = plot_df.dropna(subset=["1Y"])
    if plot_df.empty:
        return None

    plot_df = plot_df.head(8)

    fig, ax = plt.subplots(figsize=(8, 5.5))
    ax.bar(plot_df["펀드명"], plot_df["1Y"])
    ax.set_title("추천 펀드 1년 수익률")
    ax.set_xlabel("펀드")
    ax.set_ylabel("수익률(%)")
    plt.xticks(rotation=45, ha="right")

    output_path = charts_dir / "fund_return_1y.png"
    fig.savefig(output_path, bbox_inches="tight", dpi=200)
    plt.close(fig)
    return str(output_path)


def create_allocation_chart(allocation_df: pd.DataFrame, charts_dir: Path):
    if allocation_df.empty:
        return None

    fig, ax = plt.subplots(figsize=(5.5, 5.5))
    ax.pie(
        allocation_df["weight"],
        labels=allocation_df["asset_name"],
        autopct="%1.1f%%",
        startangle=90,
    )
    ax.set_title("당월 자산배분")

    output_path = charts_dir / "allocation_current.png"
    fig.savefig(output_path, bbox_inches="tight", dpi=200)
    plt.close(fig)
    return str(output_path)


def prepare_table_df(df: pd.DataFrame, keep_columns: list[str], rename_map: dict[str, str], max_rows=5):
    if df.empty:
        return pd.DataFrame()

    display_df = df.copy()
    display_df = display_df[keep_columns].head(max_rows)
    display_df = display_df.rename(columns=rename_map)

    text_columns = ["구분", "자산명", "펀드명", "자산군", "스타일", "운용사", "수탁회사", "사무관리회사", "유형", "펀드등급", "설정일"]

    for col in display_df.columns:
        if col not in text_columns:
            try:
                numeric_series = pd.to_numeric(display_df[col])
                display_df[col] = numeric_series.map(
                    lambda x: f"{x:.2f}" if pd.notnull(x) else ""
                )
            except Exception:
                display_df[col] = display_df[col].astype(str)

    return display_df


def add_table_slide(prs, report_month: str, section_no: str, title: str, df: pd.DataFrame, widths, top=2.45):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_report_header(slide, report_month, section_no, title)

    if df.empty:
        add_textbox(slide, 1.0, 2.6, 5.0, 1.0, "표 데이터가 없습니다.")
        return slide

    rows = len(df) + 1
    cols = len(df.columns)
    table = slide.shapes.add_table(
        rows,
        cols,
        Inches(0.35),
        Inches(top),
        Inches(7.45),
        Inches(3.55),
    ).table

    set_column_widths(table, widths)

    for c, col_name in enumerate(df.columns):
        table.cell(0, c).text = str(col_name)
        style_table_header(table.cell(0, c))

    for r in range(len(df)):
        for c in range(cols):
            value = str(df.iloc[r, c])
            table.cell(r + 1, c).text = value
            align = PP_ALIGN.LEFT if c in [0, 1] else PP_ALIGN.CENTER
            style_table_body(table.cell(r + 1, c), align)

    return slide


def build_summary_text(month, source_status, market_comments, analyst_text, manager_text):
    now = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    existing_count = sum(1 for item in source_status if item["exists"])
    total_count = len(source_status)

    lines = [
        "Monthly Investment Report Draft Summary",
        "=" * 40,
        f"Generated at: {now}",
        f"Report month: {month}",
        "",
        f"Source files found: {existing_count}/{total_count}",
        "",
        "[글로벌 경기 전망]",
        str(market_comments.get("global_outlook", ""))[:800],
        "",
        "[국내주식]",
        str(market_comments.get("domestic_stocks", ""))[:800],
        "",
        "[국내채권]",
        str(market_comments.get("domestic_bonds", ""))[:800],
        "",
        "[애널리스트 코멘트 요약]",
        str(analyst_text)[:600],
        "",
        "[매니저 코멘트 요약]",
        str(manager_text)[:600],
    ]
    return "\n".join(lines)


def create_pptx_report(
    report_month: str,
    paths: dict[str, Path],
    cover_info: dict,
    market_comments: dict,
    market_perf_df: pd.DataFrame,
    analyst_text: str,
    manager_text: str,
    allocation_reason: str,
    asset_chart_path: str | None,
    market_chart_1: str | None,
    market_chart_2: str | None,
    fund_chart_path: str | None,
    allocation_chart_path: str | None,
    fund_table: pd.DataFrame,
    perf_table: pd.DataFrame,
    manager_status_table: pd.DataFrame,
):
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    page_no = 1

    # 1. Cover
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_cover_background(slide)

    add_textbox(
        slide, 0.7, 1.6, 6.8, 0.8,
        cover_info.get("title", "동양생명 월간 투자전략"),
        font_size=24, bold=True, font_color=WHITE
    )
    add_textbox(
        slide, 0.7, 2.3, 6.8, 0.4,
        cover_info.get("subtitle", "변액보험 월간 투자전략 GUIDE"),
        font_size=14, font_color=WHITE
    )
    add_textbox(
        slide, 0.7, 3.1, 3.5, 0.3,
        f"기준월: {report_month}",
        font_size=13, font_color=WHITE
    )
    add_textbox(
        slide, 0.7, 3.5, 4.5, 0.3,
        f"기준일: {market_comments.get('base_date', '')}",
        font_size=13, font_color=WHITE
    )

    notice_y = 4.3
    for notice in cover_info.get("notices", [])[:3]:
        add_textbox(
            slide,
            0.8,
            notice_y,
            6.8,
            0.35,
            f"• {notice}",
            font_size=9,
            font_color=WHITE
        )
        notice_y += 0.38

    add_textbox(
        slide,
        0.5,
        11.0,
        3.8,
        0.2,
        f"동양생명 월간 투자전략 | {report_month}",
        font_size=8,
        font_color=WHITE,
    )
    add_textbox(
        slide,
        7.1,
        11.0,
        0.5,
        0.2,
        str(page_no),
        font_size=8,
        align=PP_ALIGN.RIGHT,
        font_color=WHITE,
    )
    page_no += 1

    # 2. Key summary
    perf_display = prepare_table_df(
        market_perf_df,
        keep_columns=["구분", "자산명", "1개월", "3개월", "6개월", "1년", "3년"],
        rename_map={},
        max_rows=10,
    )
    slide = add_table_slide(
        prs,
        report_month,
        "01",
        "핵심 요약",
        perf_display,
        [0.75, 1.55, 0.72, 0.72, 0.72, 0.72, 0.72],
        top=2.35
    )

    add_textbox(
        slide,
        0.55,
        6.18,
        7.1,
        0.25,
        f"기준일: {market_comments.get('base_date', '')} / 주요 자산별 성과 요약",
        font_size=8.8,
        font_color=THEME_DARK,
    )

    analyst_summary = str(analyst_text).replace("\n", " ")[:260]
    manager_summary = str(manager_text).replace("\n", " ")[:260]

    add_summary_comment_box(
        slide,
        0.50,
        6.55,
        3.45,
        2.05,
        "애널리스트 코멘트 요약",
        analyst_summary,
    )
    add_summary_comment_box(
        slide,
        4.10,
        6.55,
        3.45,
        2.05,
        "펀드매니저 코멘트 요약",
        manager_summary,
    )

    add_footer(slide, report_month, page_no)
    page_no += 1

    # 3. Major asset market inspection
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_report_header(slide, report_month, "02", "주요 자산 시장 점검")

    add_chart_comment_box(
        slide,
        0.55,
        2.30,
        7.10,
        1.45,
        "글로벌 경기 전망",
        str(market_comments.get("global_outlook", ""))[:420] or "코멘트 없음",
    )
    add_chart_comment_box(
        slide,
        0.55,
        3.95,
        7.10,
        1.45,
        "국내주식",
        str(market_comments.get("domestic_stocks", ""))[:420] or "코멘트 없음",
    )
    add_chart_comment_box(
        slide,
        0.55,
        5.60,
        7.10,
        1.45,
        "국내채권",
        str(market_comments.get("domestic_bonds", ""))[:420] or "코멘트 없음",
    )

    if market_chart_1 and Path(market_chart_1).exists():
        slide.shapes.add_picture(market_chart_1, Inches(0.55), Inches(7.35), width=Inches(3.35))
    if market_chart_2 and Path(market_chart_2).exists():
        slide.shapes.add_picture(market_chart_2, Inches(4.20), Inches(7.35), width=Inches(3.35))

    add_footer(slide, report_month, page_no)
    page_no += 1

    # 4. Analyst
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_report_header(slide, report_month, "03", "애널리스트 코멘트")
    add_profile_box(
        slide,
        left=5.45,
        top=2.30,
        width=2.15,
        height=1.65,
        company=ANALYST_PROFILE["company"],
        name=ANALYST_PROFILE["name"],
        photo_path=ANALYST_PROFILE["photo"],
    )
    add_textbox(slide, 0.6, 2.35, 4.55, 7.8, analyst_text[:2200], font_size=10.8)
    add_footer(slide, report_month, page_no)
    page_no += 1

    # 5. Manager
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_report_header(slide, report_month, "04", "매니저 코멘트")
    add_profile_box(
        slide,
        left=5.45,
        top=2.30,
        width=2.15,
        height=1.65,
        company=MANAGER_PROFILE["company"],
        name=MANAGER_PROFILE["name"],
        photo_path=MANAGER_PROFILE["photo"],
    )
    add_textbox(slide, 0.6, 2.35, 4.55, 7.8, manager_text[:2200], font_size=10.8)
    add_footer(slide, report_month, page_no)
    page_no += 1

    # 6. Allocation
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_report_header(slide, report_month, "05", "투자자별 자산배분 전략")
    add_chart_comment_box(
        slide,
        0.6,
        2.35,
        7.0,
        1.7,
        "자산배분 코멘트",
        allocation_reason[:380],
    )
    if allocation_chart_path and Path(allocation_chart_path).exists():
        slide.shapes.add_picture(allocation_chart_path, Inches(1.4), Inches(4.4), width=Inches(5.1))
    add_footer(slide, report_month, page_no)
    page_no += 1

    # 7. Fund table
    fund_display = prepare_table_df(
        fund_table,
        keep_columns=["유형", "펀드명", "펀드등급", "1Y", "3Y"],
        rename_map={},
        max_rows=5,
    )
    slide = add_table_slide(prs, report_month, "06", "추천 펀드 수익률 현황", fund_display, [1.2, 3.6, 1.2, 0.9, 0.9])
    add_footer(slide, report_month, page_no)
    page_no += 1

    # 8. Fund chart
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_report_header(slide, report_month, "07", "추천 펀드 차트")
    if fund_chart_path and Path(fund_chart_path).exists():
        slide.shapes.add_picture(fund_chart_path, Inches(0.6), Inches(2.35), width=Inches(7.0))
    add_chart_comment_box(
        slide,
        0.6,
        7.85,
        7.0,
        1.5,
        "포인트",
        "추천 펀드 간 1년 수익률 비교를 중심으로 배치했습니다.",
    )
    add_footer(slide, report_month, page_no)
    page_no += 1

    # 9. Performance
    perf_table_display = prepare_table_df(
        perf_table,
        keep_columns=["자산군", "펀드명", "1M", "3M", "1Y"],
        rename_map={},
        max_rows=5,
    )
    slide = add_table_slide(prs, report_month, "08", "변액펀드 성과현황", perf_table_display, [1.2, 3.5, 0.9, 0.9, 0.9])
    add_footer(slide, report_month, page_no)
    page_no += 1

    # 10. Manager status
    manager_display = prepare_table_df(
        manager_status_table,
        keep_columns=["펀드명", "운용사", "수탁회사"],
        rename_map={},
        max_rows=5,
    )
    slide = add_table_slide(prs, report_month, "09", "펀드별 위탁운용사 현황", manager_display, [2.4, 2.6, 2.0])
    add_footer(slide, report_month, page_no)

    output_path = paths["report_root"] / "monthly_report_draft.pptx"
    prs.save(output_path)
    return output_path


def write_outputs(paths: dict[str, Path], summary_text: str, metadata: dict, pptx_path: Path):
    summary_path = paths["report_root"] / "report_summary.txt"
    metadata_path = paths["report_root"] / "report_metadata.json"

    summary_path.write_text(summary_text, encoding="utf-8")
    metadata_path.write_text(
        json.dumps(metadata, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )

    print(f"Created: {summary_path}")
    print(f"Created: {metadata_path}")
    print(f"Created: {pptx_path}")
    print("Monthly report draft generation completed.")


def main():
    args = parse_args()
    paths = ensure_directories(args.month, args.mode)

    source_status = summarize_source_status()

    cover_info = get_cover_info()
    market_comments = get_market_comments()
    market_perf_df = get_market_perf_table()
    analyst_text = get_analyst_comment()
    manager_text = get_manager_comment()
    allocation_info = get_allocation_data()
    fund_table = get_recommended_fund_table()
    perf_table = get_performance_table()
    manager_status_table = get_manager_status_table()

    asset_chart_path = create_market_perf_summary_chart(market_perf_df, paths["charts_dir"])
    market_chart_1, market_chart_2 = create_market_dual_chart(market_perf_df, paths["charts_dir"])
    fund_chart_path = create_fund_chart(fund_table, paths["charts_dir"])
    allocation_chart_path = create_allocation_chart(allocation_info["current_alloc"], paths["charts_dir"])

    summary_text = build_summary_text(
        args.month,
        source_status,
        market_comments,
        analyst_text,
        manager_text,
    )

    metadata = {
        "generated_at": datetime.now().isoformat(),
        "report_month": args.month,
        "mode": args.mode,
        "source_status": source_status,
        "profiles": {
            "analyst_company": ANALYST_PROFILE["company"],
            "analyst_name": ANALYST_PROFILE["name"],
            "manager_company": MANAGER_PROFILE["company"],
            "manager_name": MANAGER_PROFILE["name"],
        },
        "tables": {
            "market_perf_rows": len(market_perf_df),
            "fund_rows": len(fund_table),
            "performance_rows": len(perf_table),
            "manager_status_rows": len(manager_status_table),
        },
        "comments": {
            "global_outlook_len": len(market_comments.get("global_outlook", "")),
            "domestic_stocks_len": len(market_comments.get("domestic_stocks", "")),
            "domestic_bonds_len": len(market_comments.get("domestic_bonds", "")),
        },
        "charts": {
            "asset_chart": asset_chart_path,
            "market_chart_1": market_chart_1,
            "market_chart_2": market_chart_2,
            "fund_chart": fund_chart_path,
            "allocation_chart": allocation_chart_path,
        },
    }

    pptx_path = create_pptx_report(
        report_month=args.month,
        paths=paths,
        cover_info=cover_info,
        market_comments=market_comments,
        market_perf_df=market_perf_df,
        analyst_text=analyst_text,
        manager_text=manager_text,
        allocation_reason=allocation_info.get("reason", ""),
        asset_chart_path=asset_chart_path,
        market_chart_1=market_chart_1,
        market_chart_2=market_chart_2,
        fund_chart_path=fund_chart_path,
        allocation_chart_path=allocation_chart_path,
        fund_table=fund_table,
        perf_table=perf_table,
        manager_status_table=manager_status_table,
    )

    write_outputs(paths, summary_text, metadata, pptx_path)


if __name__ == "__main__":
    main()
